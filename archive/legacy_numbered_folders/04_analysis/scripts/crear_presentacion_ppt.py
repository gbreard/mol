"""
Generador de Presentación PowerPoint
=====================================

Crea una presentación completa con todas las visualizaciones
del análisis de ofertas laborales ZonaJobs + ESCO.

Fecha: 2025-10-17
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from datetime import datetime


def crear_portada(prs):
    """Crea slide de portada."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Fondo azul
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 58, 138)  # Azul oscuro

    # Título principal
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "Análisis de Ofertas Laborales"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    top = Inches(3.5)
    subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(1))
    text_frame = subtitle_box.text_frame

    p = text_frame.paragraphs[0]
    p.text = "ZonaJobs.com.ar + Taxonomía ESCO"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER

    # Info adicional
    top = Inches(5.5)
    info_box = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    text_frame = info_box.text_frame

    lines = [
        "61 Ofertas Analizadas | 73.8% Clasificadas",
        "Período: Agosto - Octubre 2025",
        "OEDE - Observatorio de Empleo y Dinámica Empresarial"
    ]

    for i, line in enumerate(lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    return slide


def crear_slide_seccion(prs, titulo, subtitulo=None):
    """Crea slide divisor de sección."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo gradiente (simulado con color sólido)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(59, 130, 246)  # Azul medio

    # Título
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame

    p = text_frame.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    if subtitulo:
        top = Inches(4)
        subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(1))
        text_frame = subtitle_box.text_frame

        p = text_frame.paragraphs[0]
        p.text = subtitulo
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(220, 240, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide


def crear_slide_imagen(prs, titulo, imagen_path, descripcion=None):
    """Crea slide con título e imagen."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame

    p = text_frame.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    p.alignment = PP_ALIGN.CENTER

    # Imagen centrada
    img_left = Inches(0.5)
    img_top = Inches(1.2)
    img_width = Inches(9)

    try:
        slide.shapes.add_picture(str(imagen_path), img_left, img_top, width=img_width)
    except Exception as e:
        # Si falla, agregar texto de error
        error_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        p = error_box.text_frame.paragraphs[0]
        p.text = f"Error al cargar imagen: {imagen_path.name}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 0, 0)

    # Descripción (si existe)
    if descripcion:
        desc_top = Inches(6.8)
        desc_box = slide.shapes.add_textbox(left, desc_top, width, Inches(0.5))
        text_frame = desc_box.text_frame

        p = text_frame.paragraphs[0]
        p.text = descripcion
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    return slide


def crear_slide_objetivos(prs):
    """Crea slide de objetivos y alcance."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Objetivos y Alcance del Proyecto"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    p.alignment = PP_ALIGN.CENTER

    # Contenido en dos columnas
    # Columna 1: Objetivos
    left1 = Inches(0.5)
    top1 = Inches(1.2)
    width1 = Inches(4.25)

    obj_box = slide.shapes.add_textbox(left1, top1, width1, Inches(0.5))
    p = obj_box.text_frame.paragraphs[0]
    p.text = "OBJETIVOS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # Lista de objetivos
    top1 = Inches(1.8)
    obj_text_box = slide.shapes.add_textbox(left1, top1, width1, Inches(5))
    text_frame = obj_text_box.text_frame
    text_frame.word_wrap = True

    objetivos = [
        "Extraer ofertas laborales de ZonaJobs.com.ar mediante web scraping",
        "Clasificar automáticamente ofertas usando taxonomía ESCO",
        "Enriquecer ofertas con skills y competencias",
        "Analizar patrones temporales de publicación",
        "Identificar ocupaciones más demandadas",
        "Generar insights sobre el mercado laboral"
    ]

    for i, obj in enumerate(objetivos):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = "• " + obj
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.line_spacing = 1.3

    # Columna 2: Alcance
    left2 = Inches(5.25)
    alcance_box = slide.shapes.add_textbox(left2, top1 - Inches(0.6), width1, Inches(0.5))
    p = alcance_box.text_frame.paragraphs[0]
    p.text = "ALCANCE"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # Datos del alcance
    alcance_text_box = slide.shapes.add_textbox(left2, Inches(1.8), width1, Inches(5))
    text_frame = alcance_text_box.text_frame
    text_frame.word_wrap = True

    alcance = [
        "Dataset: 61 ofertas laborales",
        "Período: Agosto - Octubre 2025",
        "Fuente: ZonaJobs.com.ar",
        "Taxonomía: ESCO v1.1 (3,046 ocupaciones)",
        "Campos extraídos: 33 por oferta",
        "Clasificación: 73.8% exitosa",
        "Skills identificadas: 107 únicas",
        "Códigos ISCO: 100% cobertura"
    ]

    for i, item in enumerate(alcance):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = "✓ " + item
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.line_spacing = 1.3

    return slide


def crear_slide_metodologia(prs):
    """Crea slide de metodología técnica."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Metodología del Proceso"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    p.alignment = PP_ALIGN.CENTER

    # Flujo del proceso
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)

    text_box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    pasos = [
        ("1. WEB SCRAPING", [
            "Descubrimiento de API de ZonaJobs mediante análisis de tráfico de red",
            "Extracción automatizada vía requests HTTP con rate limiting (2s)",
            "Parser de 33 campos por oferta (título, empresa, modalidad, fechas, etc.)"
        ]),
        ("2. PROCESAMIENTO Y LIMPIEZA", [
            "Normalización de textos (acentos, mayúsculas, caracteres especiales)",
            "Eliminación de HTML y formateo",
            "Estandarización de fechas y categorías"
        ]),
        ("3. CLASIFICACIÓN SEMÁNTICA (ESCO)", [
            "Matching semántico basado en similitud de coseno (TF-IDF)",
            "Normalización de vocabulario español España ↔ Argentina",
            "Threshold de similitud: 0.4 (balanceando precisión vs cobertura)"
        ]),
        ("4. ENRIQUECIMIENTO", [
            "Asignación de códigos ISCO (1D, 2D, 3D, 4D) desde RDF",
            "Extracción de skills esenciales y opcionales por ocupación",
            "Vinculación con 6,818 skills de la ontología ESCO"
        ]),
        ("5. ANÁLISIS Y VISUALIZACIÓN", [
            "Análisis estadístico multidimensional (ocupaciones, skills, temporal)",
            "Generación de 13 visualizaciones estáticas (300 DPI)",
            "Dashboard interactivo con Plotly"
        ])
    ]

    for i, (titulo, items) in enumerate(pasos):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[sum(len(items) + 1 for _, items in pasos[:i])]
        p.text = titulo
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = RGBColor(59, 130, 246)
        p.space_before = Pt(10) if i > 0 else Pt(0)

        for j, item in enumerate(items):
            text_frame.add_paragraph()
            p = text_frame.paragraphs[-1]
            p.text = "  • " + item
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(2)
            p.line_spacing = 1.2

    return slide


def crear_slide_herramientas(prs):
    """Crea slide de herramientas y tecnologías."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Stack Tecnológico"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    p.alignment = PP_ALIGN.CENTER

    # Tres columnas
    col_width = Inches(2.8)
    col_height = Inches(5.5)
    top_start = Inches(1.3)

    columnas = [
        ("SCRAPING & DATOS", [
            "Python 3.13",
            "Requests (HTTP)",
            "Pandas (procesamiento)",
            "JSON / CSV / Excel",
            "Playwright (opcional)"
        ]),
        ("ANÁLISIS SEMÁNTICO", [
            "Scikit-learn (TF-IDF)",
            "ESCO v1.1 (RDF)",
            "RDFLib (parsing)",
            "NumPy",
            "Normalización ES-AR"
        ]),
        ("VISUALIZACIÓN", [
            "Matplotlib",
            "Plotly (dashboards)",
            "Python-pptx",
            "Seaborn",
            "HTML/CSS"
        ])
    ]

    for i, (titulo, items) in enumerate(columnas):
        left_col = Inches(0.7 + i * 3.1)

        # Título columna
        title_col_box = slide.shapes.add_textbox(left_col, top_start, col_width, Inches(0.5))
        p = title_col_box.text_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(59, 130, 246)
        p.alignment = PP_ALIGN.CENTER

        # Items
        items_box = slide.shapes.add_textbox(left_col, top_start + Inches(0.6), col_width, col_height - Inches(0.6))
        text_frame = items_box.text_frame
        text_frame.word_wrap = True

        for j, item in enumerate(items):
            if j > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[j]
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(8)
            p.line_spacing = 1.3

    # Nota al pie
    nota_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    p = nota_box.text_frame.paragraphs[0]
    p.text = "Código completo disponible en: D:\\OEDE\\Webscrapping"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.CENTER

    return slide


def crear_slide_conclusiones(prs):
    """Crea slide de conclusiones."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)

    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Conclusiones Principales"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    p.alignment = PP_ALIGN.CENTER

    # Conclusiones en bullets
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    conclusiones = [
        "73.8% de ofertas clasificadas exitosamente con taxonomía ESCO",
        "Predominio de profesionales científicos (26.7%) y técnicos (22.2%)",
        "Miércoles es el día más activo para publicación de ofertas (33%)",
        "Septiembre mostró el mayor volumen de ofertas (47.5% del total)",
        "100% de ofertas clasificadas tienen código ISCO asignado",
        "Similitud promedio de matching: 0.537 (con 2 matches perfectos)",
        "Mayor demanda de skills administrativo-contables y sanitarias",
        "Patrón claro: publicación concentrada en primera mitad de semana"
    ]

    for i, conclusion in enumerate(conclusiones):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = conclusion
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.level = 0
        p.space_before = Pt(8)
        p.line_spacing = 1.2

    return slide


def crear_slide_cierre(prs):
    """Crea slide de cierre."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 58, 138)

    # Mensaje de cierre
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    lines = [
        "¡Gracias!",
        "",
        "Proyecto: Web Scraping ZonaJobs + ESCO",
        "OEDE - Observatorio de Empleo y Dinámica Empresarial",
        f"Fecha: {datetime.now().strftime('%B %Y')}"
    ]

    for i, line in enumerate(lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        if i == 0:
            p.font.size = Pt(48)
            p.font.bold = True
        else:
            p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    return slide


def main():
    """Función principal."""
    print("="*70)
    print("GENERADOR DE PRESENTACION POWERPOINT")
    print("Analisis ZonaJobs + ESCO")
    print("="*70)

    # Rutas
    base_dir = Path(r"D:\OEDE\Webscrapping")
    charts_dir = base_dir / "data" / "processed" / "charts"
    output_dir = base_dir / "data" / "processed"

    # Verificar directorio
    if not charts_dir.exists():
        print(f"\n[ERROR] No se encontro el directorio: {charts_dir}")
        return

    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\n>> Creando presentacion...")

    # 1. Portada
    print("   [+] Portada")
    crear_portada(prs)

    # 2. Metodología y Contexto
    print("   [+] Objetivos y Alcance")
    crear_slide_objetivos(prs)

    print("   [+] Metodologia del Proceso")
    crear_slide_metodologia(prs)

    print("   [+] Stack Tecnologico")
    crear_slide_herramientas(prs)

    # 3. Sección: Clasificación ESCO
    print("   [+] Seccion: Clasificacion ESCO")
    crear_slide_seccion(prs, "Clasificación con Taxonomía ESCO",
                        "Análisis de ocupaciones, grupos ISCO y calidad de matching")

    # Slides de clasificación ESCO
    imagenes_esco = [
        ("01_top_ocupaciones.png", "Top 15 Ocupaciones ESCO Identificadas",
         "37 ocupaciones únicas identificadas en el dataset"),
        ("02_distribucion_isco.png", "Distribución por Grupos ISCO (2 dígitos)",
         "8 grupos principales representados"),
        ("07_isco_grupos_principales.png", "Distribución por Grupos ISCO (1 dígito)",
         "Predominio de profesionales científicos y técnicos"),
        ("06_tasa_clasificacion.png", "Tasa de Clasificación Exitosa",
         "73.8% de ofertas clasificadas con ESCO"),
        ("03_distribucion_similitud.png", "Calidad del Matching Semántico",
         "Similitud promedio: 0.537, con 2 matches perfectos"),
    ]

    for img_name, titulo, desc in imagenes_esco:
        img_path = charts_dir / img_name
        if img_path.exists():
            print(f"   [+] {titulo}")
            crear_slide_imagen(prs, titulo, img_path, desc)
        else:
            print(f"   [!] No encontrado: {img_name}")

    # 4. Sección: Skills y Competencias
    print("   [+] Seccion: Skills y Competencias")
    crear_slide_seccion(prs, "Skills y Competencias",
                        "Análisis de habilidades demandadas en las ofertas")

    imagenes_skills = [
        ("04_top_skills.png", "Top 15 Skills Más Demandadas",
         "82.2% de ofertas enriquecidas con skills esenciales"),
        ("08_skills_por_oferta.png", "Distribución de Skills por Oferta",
         "Promedio: 2.9 skills esenciales por oferta"),
    ]

    for img_name, titulo, desc in imagenes_skills:
        img_path = charts_dir / img_name
        if img_path.exists():
            print(f"   [+] {titulo}")
            crear_slide_imagen(prs, titulo, img_path, desc)
        else:
            print(f"   [!] No encontrado: {img_name}")

    # 5. Sección: Modalidad y Geografía
    print("   [+] Seccion: Modalidad y Geografia")
    crear_slide_seccion(prs, "Modalidad de Trabajo y Distribución",
                        "Análisis de modalidades laborales y ubicación geográfica")

    imagenes_modalidad = [
        ("05_modalidad_trabajo.png", "Distribución por Modalidad de Trabajo",
         "78.7% presencial, 21.3% híbrido"),
    ]

    # Verificar si existen imágenes adicionales
    if (charts_dir / "09_distribucion_geografica.png").exists():
        imagenes_modalidad.append(
            ("09_distribucion_geografica.png", "Distribución Geográfica de Ofertas",
             "Análisis por ubicación de las ofertas")
        )

    if (charts_dir / "10_top_empresas.png").exists():
        imagenes_modalidad.append(
            ("10_top_empresas.png", "Top Empresas con Más Ofertas",
             "Empresas más activas en publicación de ofertas")
        )

    for img_name, titulo, desc in imagenes_modalidad:
        img_path = charts_dir / img_name
        if img_path.exists():
            print(f"   [+] {titulo}")
            crear_slide_imagen(prs, titulo, img_path, desc)
        else:
            print(f"   [!] No encontrado: {img_name}")

    # 6. Sección: Análisis Temporal
    print("   [+] Seccion: Analisis Temporal")
    crear_slide_seccion(prs, "Análisis Temporal",
                        "Evolución de ofertas en el tiempo (Agosto - Octubre 2025)")

    imagenes_temporal = [
        ("09_serie_temporal_diaria.png", "Serie Temporal Diaria de Ofertas",
         "Promedio: 1.91 ofertas/día, Máximo: 4 ofertas"),
        ("10_serie_temporal_semanal.png", "Serie Temporal Semanal de Ofertas",
         "Promedio: 6.78 ofertas/semana, Pico en semana 34 (12 ofertas)"),
        ("11_serie_diaria_isco_stacked.png", "Serie Diaria por Grupo Ocupacional ISCO",
         "Composición ocupacional estable a lo largo del tiempo"),
        ("12_serie_semanal_isco_stacked.png", "Serie Semanal por Grupo Ocupacional ISCO",
         "Mayor diversidad en semanas de alto volumen"),
        ("13_heatmap_dia_mes.png", "Distribución por Día de Semana y Mes",
         "Miércoles: día más activo (33% de ofertas)"),
    ]

    for img_name, titulo, desc in imagenes_temporal:
        img_path = charts_dir / img_name
        if img_path.exists():
            print(f"   [+] {titulo}")
            crear_slide_imagen(prs, titulo, img_path, desc)
        else:
            print(f"   [!] No encontrado: {img_name}")

    # 7. Conclusiones
    print("   [+] Conclusiones")
    crear_slide_conclusiones(prs)

    # 8. Cierre
    print("   [+] Cierre")
    crear_slide_cierre(prs)

    # Guardar presentación
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = output_dir / f"Presentacion_ZonaJobs_ESCO_{timestamp}.pptx"

    prs.save(str(output_file))

    print("\n" + "="*70)
    print("PRESENTACION GENERADA EXITOSAMENTE")
    print("="*70)
    print(f"\n>> Archivo: {output_file.name}")
    print(f">> Ubicacion: {output_file.parent}")
    print(f">> Slides totales: {len(prs.slides)}")
    print("\n[OK] Presentacion lista para usar!")
    print("="*70)

    return output_file


if __name__ == "__main__":
    main()
